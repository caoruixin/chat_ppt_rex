from pptx import Presentation

def remove_all_slides(prs: Presentation):
    # xml_slides = prs.slides._sldIdLst
    # slides = list(xml_slides)
    # for slide in slides:
    #     xml_slides.remove(slide)
    # print("所有默认幻灯片已被移除。")

    """Remove all slides from a PowerPoint presentation."""
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids:
        rId = slide_id.rId
        prs.part.drop_rel(rId)  # 删除关联关系
        prs.slides._sldIdLst.remove(slide_id)  # 从幻灯片列表中移除
